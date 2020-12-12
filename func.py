
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import copy
import six
import os
import sys
from xml.etree import ElementTree
import xml.etree.ElementTree as ET
import xml.dom.minidom
import os
import slot_adder, ratio_comparision


a_float = 3.14159
formatted_float = "{:.2f}".format(a_float)
# Format the float with two decimal places

# may be added for the certainity of ratios because they vary in calculations
print(formatted_float)


# RATİOS WİLL BE TAKEN ACCORDİNG TO WİDTH/HEİGHT

card_details = [{
    'name': 'card-1',
    'width': 100,
    'height': 300
}, {
    'name': 'card-2',
    'width': 200,
    'height': 200
}, {
    'name': 'card-3',
    'width': 300,
    'height': 200
}, {
    'name': 'card-4',
    'width': 600,   # even 400 t0 200 gave column, maybe added formatter ?

    'height': 200
}]


_orderofnums = [0, 1, 2, 3, 4, 5]


class KeyFound(Exception):
    def __init__(self, message):
        self.value = message


def XML(shape, key, value=None):
    try:
        XMLTraversal(shape.element, key, value)
    except KeyFound as e:
        return e.value
    except:
        pass
    return None


def XMLTraversal(shape_element, key, value=None):
    for child in shape_element.getchildren():
        if child.attrib:
            if key in child.attrib:
                if value:
                    child.attrib[key] = value
                raise KeyFound(child.attrib[key])
        XMLTraversal(child, key, value)


def ratio_exe(column, bar, car, slot_num, halfbar=None, square=None):
    column_ratio = column
    bar_ratio = bar
    ratio_car = car["width"] / car["height"]
    diffcolcar = column_ratio - ratio_car
    diffbarcar = bar_ratio - ratio_car

    list_dict = [
        {"name": "column", "diff": abs(diffcolcar)},
        {"name": "bar", "diff": abs(diffbarcar)}

    ]

    print("********************************")
    half_bar = halfbar
    diffhalfbarcar = half_bar - ratio_car
    list_dict.append({"name": "half_bar", "diff": abs(diffhalfbarcar)})

    square_ratio = square
    square_diff = square_ratio - ratio_car
    list_dict.append({"name": "square", "diff": abs(square_diff)})

    list_diff = []
    for i in list_dict:
        list_diff.append(i["diff"])

    for ratio in list_diff:
        if min(list_diff) == ratio:
            for record in list_dict:
                if record["diff"] == ratio:
                    # denenebılır daha dınamık olması acısından
                    print(record["name"]+" is more suitable where " +
                          str(slot_num)+" added for "+car["name"])


def _get_blank_slide_layout(pres, layout_num):
    # which layout you want
    return pres.slide_layouts[layout_num]


def duplicate_slide(pres, index):

    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres, 2)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return dest


def add_slots(i, pres, list_of_ratios, slots_func):
    for shape in pres.slides[i].shapes:
        alt = XML(shape, "descr")
        if alt == "type:content-area":
            ratio_of_slot = slots_func(pres.slides[i], shape)
            list_of_ratios.append(ratio_of_slot)


def provider(i, pres, list_of_ratios):

    if i == 0:  # one full
        add_slots(i, pres, list_of_ratios, slot_adder.one_full)
    elif i == 1:  # two_column

        add_slots(i, pres, list_of_ratios, slot_adder.two_column)
    elif i == 2:  # two_bar
        add_slots(i, pres, list_of_ratios, slot_adder.two_bar)
    elif i == 3:  # three_column
        add_slots(i, pres, list_of_ratios, slot_adder.three_column)
    elif i == 4:  # three_bar
        add_slots(i, pres, list_of_ratios, slot_adder.three_bar)
    elif i == 5:  # four_column
        add_slots(i, pres, list_of_ratios, slot_adder.four_column)
    elif i == 6:  # four_bar
        add_slots(i, pres, list_of_ratios, slot_adder.four_bar)
    elif i == 7:  # four_half_bar
        add_slots(i, pres, list_of_ratios, slot_adder.four_half_bar)
    elif i == 8:  # five_column
        add_slots(i, pres, list_of_ratios, slot_adder.five_column)
    elif i == 9:  # five_bar
        add_slots(i, pres, list_of_ratios, slot_adder.five_bar)
    elif i == 10:  # five_half_bar
        add_slots(i, pres, list_of_ratios, slot_adder.five_half_bar)
    elif i == 11:  # five_square
        add_slots(i, pres, list_of_ratios, slot_adder.five_square)
    elif i == 12:  # six_square
        add_slots(i, pres, list_of_ratios, slot_adder.six_square)
    elif i == 13:  # six_column
        add_slots(i, pres, list_of_ratios, slot_adder.six_column)
    elif i == 14:  # six_half_bar
        add_slots(i, pres, list_of_ratios, slot_adder.six_half_bar)
    else:
        print("something went wrong when applying slots to slides")


content_slides = Presentation("layout-new.pptx")


def delete_slide(prs, slide):  # for deleting the last blank slide
  
    id_dict = {slide.id: [i, slide.rId]
               for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def multiply_slide(pres, times, card_details): # multiple slides for given times
    list_of_ratios = []
    for i in range(times):
        a = duplicate_slide(pres, 0), # duplicate slides

        provider(i, pres, list_of_ratios) # adding slots on duplicated slides

    delete_slide(pres, a[-1])  # delete last blank slide

    for car in card_details:

        for i in _orderofnums:
            ratio_comparision.ratio_comparison(list_of_ratios,i,car,ratio_exe)


 
 


multiply_slide(content_slides, 15, card_details)


content_slides.save("export3.pptx")
os.startfile("export3.pptx")
