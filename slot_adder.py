from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import copy
import six
import os
import sys
from xml.etree import ElementTree
import xml.etree.ElementTree as ET
import xml.dom.minidom

 


# adds new slots to the slide 
def one_full( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height
        slot1_width = content_area_width
        slot1_left = content_area_left 
        slot1_top = content_area_top 

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)
def two_column( slide,content_area):


        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height


        slot1_height = content_area_height
        slot1_width = content_area_width/2
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height 
        slot2_width = content_area_width/2
        slot2_top = content_area_top 
        slot2_left = slot1_left+slot1_width

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        a = slot1_width/slot1_height
        return a 
def two_bar( slide,content_area):
        content_area_height = content_area.height 
        content_area_width = content_area.width 
        content_area_left = content_area.left 
        content_area_top = content_area.top 

        slot1_height = content_area_height/2 
        slot1_width = content_area_width
        slot1_left = content_area_left 
        slot1_top = content_area_top 

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height/2
        slot2_width = content_area_width
        slot2_top = slot1_top + slot1_height
        slot2_left = content_area_left 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        a = slot1_width/slot1_height
        return a
def three_column( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height
        slot1_width = content_area_width/3
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height 
        slot2_width = content_area_width/3
        slot2_top = content_area_top 
        slot2_left = slot1_left+slot1_width

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height 
        slot3_width = content_area_width/3
        slot3_top = content_area_top 
        slot3_left= slot2_left+slot2_width

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)
        a = slot1_width/slot1_height
        return a
def three_bar(  slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height/3
        slot1_width = content_area_width
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height/3 
        slot2_width = content_area_width
        slot2_top = slot1_top+slot1_height 
        slot2_left = content_area_left 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height/3
        slot3_width = content_area_width
        slot3_top = slot2_top+slot2_height
        slot3_left= content_area_left

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)
        a = slot1_width/slot1_height
        return a
def four_column(  slide,content_area):

        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height
        slot1_width = content_area_width/4
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height  
        slot2_width = content_area_width/4
        slot2_top = slot1_top  
        slot2_left =  slot1_left+slot1_width 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height 
        slot3_width = content_area_width/4
        slot3_top = content_area_top
        slot3_left= slot2_left+slot2_width

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height 
        slot4_width = content_area_width /4
        slot4_top = content_area_top 
        slot4_left= slot3_left+slot3_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)
        a = slot1_width/slot1_height
        return a
def four_bar( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height/4
        slot1_width = content_area_width
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height/4  
        slot2_width = content_area_width
        slot2_top = slot1_top+slot1_height  
        slot2_left =  content_area_left 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height/4 
        slot3_width = content_area_width
        slot3_top = slot2_top+slot2_height
        slot3_left= content_area_left

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height /4
        slot4_width = content_area_width  
        slot4_top = slot3_top+slot3_height 
        slot4_left= content_area_left

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)
        a = slot1_width/slot1_height
        return a
def four_half_bar( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height/2
        slot1_width = content_area_width/2
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height /2  
        slot2_width = content_area_width/2
        slot2_top = content_area_top   
        slot2_left =  slot1_left+slot1_width 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height/2 
        slot3_width = content_area_width/2
        slot3_top = slot2_top+slot2_height
        slot3_left= content_area_left

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height/2
        slot4_width = content_area_width /2 
        slot4_top = slot2_top+slot3_height 
        slot4_left= slot3_left+slot3_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)
        a = slot1_width/slot1_height
        return a
def five_bar( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height/5
        slot1_width = content_area_width 
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height /5  
        slot2_width = content_area_width 
        slot2_top = slot1_top+slot1_height   
        slot2_left =  content_area_left

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height/5 
        slot3_width = content_area_width 
        slot3_top = slot2_top+slot2_height
        slot3_left= content_area_left

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height/5
        slot4_width = content_area_width   
        slot4_top = slot3_top+slot3_height 
        slot4_left= content_area_left

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)
        
        slot5_height = content_area_height/5
        slot5_width = content_area_width   
        slot5_top = slot4_top+slot4_height 
        slot5_left= content_area_left

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot5_left,slot5_top,slot5_width,slot5_height)
       
        a = slot1_width/slot1_height
        return a
def five_column( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height
        slot1_width = content_area_width/5 
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height   
        slot2_width = content_area_width /5
        slot2_top = content_area_top   
        slot2_left =  slot1_left+slot1_width

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height  
        slot3_width = content_area_width /5
        slot3_top = content_area_top 
        slot3_left= slot2_left+slot2_width

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height 
        slot4_width = content_area_width /5  
        slot4_top = content_area_top 
        slot4_left= slot3_left+slot3_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)
        
        slot5_height = content_area_height 
        slot5_width = content_area_width/5   
        slot5_top = content_area_top  
        slot5_left= slot4_left+slot4_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot5_left,slot5_top,slot5_width,slot5_height)
       
        a = slot1_width/slot1_height
        return a
def five_half_bar( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height/3
        slot1_width = content_area_width/2
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height /3  
        slot2_width = content_area_width/2
        slot2_top = content_area_top   
        slot2_left =  slot1_left+slot1_width 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height/3 
        slot3_width = content_area_width/2
        slot3_top = slot2_top+slot2_height
        slot3_left= content_area_left

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height/3
        slot4_width = content_area_width /2 
        slot4_top = slot2_top+slot3_height 
        slot4_left= slot3_left+slot3_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)

        slot5_height = slot4_height
        slot5_width = slot4_width
        slot5_left = content_area_left+(content_area_width-slot5_width)/2
        slot5_top= slot4_top+slot4_height

        Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot5_left,slot5_top,slot5_width,slot5_height)

        a = slot1_width/slot1_height
        return a
def five_square(  slide,content_area):# sss
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        if content_area_width > content_area_height: # bar 
                slot1_width = content_area_width/3
                slot1_height = content_area_height/2
                slot1_left =  content_area_left 
                slot1_top = content_area_top
                

                First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

                slot2_width = content_area_width/3
                slot2_height = content_area_height/2
                slot2_top = content_area_top   
                slot2_left =  slot1_left+slot1_width 

                Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

                slot3_width = content_area_width/3
                slot3_height = content_area_height/2
                slot3_top = content_area_top 
                slot3_left= slot2_left+slot2_width

                Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

                slot4_width = content_area_width/3
                slot4_height = content_area_height/2
                slot4_top = slot2_top+slot3_height 
                slot4_left= content_area_left+(content_area_width-(slot4_width*2))/2

                Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)

                slot5_width = content_area_width/3
                slot5_height = content_area_height/2
                slot5_left = slot4_left+slot4_width
                slot5_top= slot3_top+slot3_height

                Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot5_left,slot5_top,slot5_width,slot5_height)

                slot1_width = content_area_width/3
                slot1_height = content_area_height/2
                slot6_left = slot5_left+slot5_width
                slot6_top= slot3_top+slot3_height

                Sixth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot6_left,slot6_top,slot6_width,slot6_height)
                a = slot1_width/slot1_height
                return a
        else : # column and square 
                slot1_width = content_area_width/2
                slot1_height = slot1_width
                slot1_left =  content_area_left
                slot1_top = content_area_top

                First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

                slot2_width = content_area_width/2
                slot2_height = slot2_width
                slot2_top = content_area_top   
                slot2_left =  slot1_left+slot1_width 

                Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

                slot3_width = content_area_width/2
                slot3_height = slot3_width
                slot3_top = slot1_top+slot1_height 
                slot3_left= content_area_left

                Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

                slot4_width = content_area_width/2
                slot4_height = slot4_width
                slot4_top = slot3_top  
                slot4_left= slot3_left+slot3_width

                Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)

                slot5_width = content_area_width/2
                slot5_height = slot5_width
                slot5_left = content_area_left+(content_area_width-slot5_width)/2
                slot5_top= slot4_top+slot4_height

                Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot5_left,slot5_top,slot5_width,slot5_height)

                a = slot1_width/slot1_height
                return a
def six_square( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height
        if content_area_width > content_area_height:
                
                slot1_width = content_area_width/3
                slot1_height = content_area_height/2
                slot1_left =  content_area_left 
                slot1_top = content_area_top
                

                First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

                slot2_width = content_area_width/3
                slot2_height = content_area_height/2
                slot2_top = content_area_top   
                slot2_left =  slot1_left+slot1_width 

                Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

                slot3_width = content_area_width/3
                slot3_height = content_area_height/2
                slot3_top = content_area_top 
                slot3_left= slot2_left+slot2_width

                Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

                slot4_width = content_area_width/3
                slot4_height = content_area_height/2
                slot4_top = slot2_top+slot3_height 
                slot4_left= content_area_left

                Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)

                slot5_width = content_area_width/3
                slot5_height = content_area_height/2
                slot5_left = slot4_left+slot4_width
                slot5_top= slot3_top+slot3_height

                Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot5_left,slot5_top,slot5_width,slot5_height)

                slot1_width = content_area_width/3
                slot1_height = content_area_height/2
                slot6_left = slot5_left+slot5_width
                slot6_top= slot3_top+slot3_height

                Sixth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot6_left,slot6_top,slot6_width,slot6_height)
                a = slot1_width/slot1_height
                return a
        elif content_area_height > content_area_width: 
                 
                        
                slot1_width = content_area_width/2
                slot1_height = slot1_width
                slot1_left = content_area_left 
                slot1_top = content_area_top
                

                First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

                slot2_width = content_area_width/2
                slot2_height = slot2_width
                slot2_top = content_area_top   
                slot2_left =  slot1_left+slot1_width 
                

                Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

                slot3_width = content_area_width/2
                slot3_height = slot3_width
                slot3_top = slot1_top+slot1_height 
                slot3_left= content_area_left
                

                Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

                slot4_width = content_area_width/2
                slot4_height = slot4_width
                slot4_top = slot2_top + slot3_height
                slot4_left= slot3_left+slot3_width
                

                Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)

                slot5_width = content_area_width/2
                slot5_height = slot5_width
                slot5_left = content_area_left
                slot5_top= slot3_top+slot3_height
                

                Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot5_left,slot5_top,slot5_width,slot5_height)

                slot6_width = content_area_width/2
                slot6_height = slot6_width
                slot6_left = slot5_left+slot5_width
                slot6_top= slot3_top+slot3_height
                

                Sixth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot6_left,slot6_top,slot6_width,slot6_height)
                a = slot1_width/slot1_height
                return a
def six_column( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height
        slot1_width = content_area_width/6 
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height   
        slot2_width = content_area_width /6
        slot2_top = content_area_top   
        slot2_left =  slot1_left+slot1_width

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height  
        slot3_width = content_area_width /6
        slot3_top = content_area_top 
        slot3_left= slot2_left+slot2_width

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height 
        slot4_width = content_area_width /6  
        slot4_top = content_area_top 
        slot4_left= slot3_left+slot3_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)
        
        slot5_height = content_area_height 
        slot5_width = content_area_width/6   
        slot5_top = content_area_top  
        slot5_left= slot4_left+slot4_width

        Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot5_left,slot5_top,slot5_width,slot5_height)

        slot6_height = content_area_height 
        slot6_width = content_area_width/6
        slot6_top = content_area_top 
        slot6_left = slot5_left+slot5_width 

        Sixth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot6_left,slot6_top,slot6_width,slot6_height)
        a = slot1_width/slot1_height
        return a
def six_half_bar( slide,content_area):
        content_area_top = content_area.top 
        content_area_left = content_area.left 
        content_area_width = content_area.width 
        content_area_height = content_area.height

        slot1_height = content_area_height/3
        slot1_width = content_area_width/2
        slot1_left = content_area_left 
        slot1_top = content_area_top

        First_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot1_left,slot1_top,slot1_width,slot1_height)

        slot2_height = content_area_height /3  
        slot2_width = content_area_width/2
        slot2_top = content_area_top   
        slot2_left =  slot1_left+slot1_width 

        Second_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot2_left,slot2_top,slot2_width,slot2_height)

        slot3_height = content_area_height/3 
        slot3_width = content_area_width/2
        slot3_top = slot2_top+slot2_height
        slot3_left= content_area_left

        Third_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot3_left,slot3_top,slot3_width,slot3_height)

        slot4_height = content_area_height/3
        slot4_width = content_area_width /2 
        slot4_top = slot2_top+slot3_height 
        slot4_left= slot3_left+slot3_width

        Fourth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,slot4_left,slot4_top,slot4_width,slot4_height)

        slot5_height = slot4_height
        slot5_width = slot4_width
        slot5_left =  slot3_left 
        slot5_top= slot4_top+slot4_height

        Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot5_left,slot5_top,slot5_width,slot5_height)

        slot6_height = slot5_height
        slot6_width = slot5_width
        slot6_left =  slot4_left 
        slot6_top= slot4_top+slot4_height

        Fifth_slot = slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, slot6_left,slot6_top,slot6_width,slot6_height)

        a = slot1_width/slot1_height
        return a